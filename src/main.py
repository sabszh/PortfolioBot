# main.py

import os
import sys
from openai import OpenAI
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import docx
import csv
from bs4 import BeautifulSoup
from pptx import Presentation
import openpyxl

# Load environment variables
load_dotenv()

# Initialize OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# System prompt for the chatbot
def load_prompt():
    base_path = os.path.dirname(sys.argv[0])  # directory of the running script
    prompt_path = os.path.join(base_path,"..","prompt.txt")
    with open(prompt_path, "r", encoding="utf-8") as f:
        return f.read()

SYSTEM_PROMPT = load_prompt()

def extract_text(file):
    """Extract text from uploaded files based on their type."""
    file_type = file.name.split(".")[-1].lower()

    if file_type == "pdf":
        reader = PdfReader(file)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    elif file_type == "docx":
        doc = docx.Document(file)
        return "\n".join(p.text for p in doc.paragraphs)

    elif file_type == "txt":
        return file.read().decode("utf-8")

    elif file_type == "csv":
        content = file.read().decode("utf-8").splitlines()
        return "\n".join([", ".join(row) for row in csv.reader(content)])

    elif file_type in ["html", "htm"]:
        soup = BeautifulSoup(file.read(), "html.parser")
        return soup.get_text(separator="\n")

    elif file_type == "md":
        return file.read().decode("utf-8")

    elif file_type == "pptx":
        prs = Presentation(file)
        slides_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slides_text.append(shape.text)
        return "\n".join(slides_text)

    elif file_type == "xlsx":
        wb = openpyxl.load_workbook(file, read_only=True, data_only=True)
        text_chunks = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text_chunks.append(", ".join(str(cell) for cell in row if cell is not None))
        return "\n".join(text_chunks)

    else:
        return f"[Filformatet '{file_type}' understøttes ikke.]"

def generate_response(messages):
    """Generate a response from the OpenAI API based on the conversation history."""
    response = client.responses.create(
        model="gpt-4o",
        input=messages,
        text={"format": {"type": "text"}},
        reasoning={},
        tools=[
            {
                "type": "file_search",
                "vector_store_ids": ["vs_67efcf44a6008191ac41404094d0a98c"]
            }
        ],
        temperature=1,
        max_output_tokens=2048,
        top_p=1,
        store=True
    )
    return response.output_text
